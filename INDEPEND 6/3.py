from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]
title_1.text = "Georgian Cheese"
subtitle_1.text = "A Culinary Journey through Georgia\n[Your Name]"

# Slide 2: Introduction to Georgian Cheese
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title_2 = slide_2.shapes.title
content_2 = slide_2.placeholders[1]
title_2.text = "Introduction"
content_2.text = "Brief overview of Georgian cuisine\nImportance of cheese in Georgian diet"

# Slide 3: History of Georgian Cheese
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title_3 = slide_3.shapes.title
content_3 = slide_3.placeholders[1]
title_3.text = "Historical Background"
content_3.text = "Origin and history of cheese-making in Georgia\nAncient traditions and techniques"

# Slide 4: Popular Georgian Cheeses
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title_4 = slide_4.shapes.title
content_4 = slide_4.placeholders[1]
title_4.text = "Varieties of Georgian Cheese"
content_4.text = "Sulguni: Description and characteristics\nImeruli: Description and characteristics\nGuda: Description and characteristics\nTenili: Description and characteristics"

# Slide 5: Cheese-Making Process
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title_5 = slide_5.shapes.title
content_5 = slide_5.placeholders[1]
title_5.text = "Traditional Cheese-Making"
content_5.text = "Overview of traditional cheese-making methods\nSteps involved in making popular cheeses"

# Slide 6: Regional Specialties
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title_6 = slide_6.shapes.title
content_6 = slide_6.placeholders[1]
title_6.text = "Regional Specialties"
content_6.text = "Cheese varieties unique to specific regions in Georgia\nCultural significance and local traditions"

# Slide 7: Cheese in Georgian Cuisine
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title_7 = slide_7.shapes.title
content_7 = slide_7.placeholders[1]
title_7.text = "Culinary Uses"
content_7.text = "Traditional Georgian dishes featuring cheese\nKhachapuri: Georgia's famous cheese bread\nOther popular recipes"

# Slide 8: Cheese Festivals and Events
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title_8 = slide_8.shapes.title
content_8 = slide_8.placeholders[1]
title_8.text = "Cheese Festivals"
content_8.text = "Annual cheese festivals in Georgia\nCultural events celebrating cheese"

# Slide 9: Tasting Georgian Cheese
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title_9 = slide_9.shapes.title
content_9 = slide_9.placeholders[1]
title_9.text = "Tasting Experience"
content_9.text = "How to properly taste and enjoy Georgian cheese\nPairing cheese with wine and other foods"

# Slide 10: Conclusion
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])
title_10 = slide_10.shapes.title
content_10 = slide_10.placeholders[1]
title_10.text = "Conclusion"

# Save the presentation
file_path = "/mnt/data/Georgian_Cheese_Presentation.pptx"
prs.save(file_path)

file_path